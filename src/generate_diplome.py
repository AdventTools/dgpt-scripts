"""
generate_diplome.py

Genereaza trei PPTX-uri cu diplome individuale, din templates/model_diploma.pptx:
- output/Diplome_medici.pptx   : 1 slide per medic    (din data/medici.xlsx, dedup dupa nume)
- output/Diplome_asistenti.pptx: 1 slide per asistent (din data/asistenti.xlsx)
- output/Diplome_diverse.pptx  : 1 slide per persoana (din data/alte_diplome.xlsx, format "Titulatura Nume")

Diferenta:
- medicii au titulatura (Dr.) deja inclusa in coloana 'nume'
- asistentele nu au titulatura
- diversii primesc "{titulatura} {nume}" concatenat (ex: "Primar Exemplu Unu")

Restul template-ului e identic. Pe fiecare slide se actualizeaza:
- Shape 0: numele persoanei (font Calibri auto-scalat sa incapa pe un rand)
- Shape 2: perioada si orasul evenimentului (din event.xlsx)
"""
import openpyxl
from lxml import etree
from pptx import Presentation

from pptx_utils import (
    ROOT, DATA, TEMPLATES_DIR, OUTPUT_DIR,
    NS_A,
    CALIBRI_PATH, auto_font_pt,
    clone_slide, set_paragraph_single_run,
)

# ─────────────────────────────────────────────────────────────────
# CAI + EVENIMENT
# ─────────────────────────────────────────────────────────────────
MEDICI_XLSX    = DATA          / 'medici.xlsx'
ASISTENTI_XLSX = DATA          / 'asistenti.xlsx'
DIVERSE_XLSX   = DATA          / 'alte_diplome.xlsx'
EVENT_XLSX     = DATA          / 'event_detail.xlsx'
TEMPLATE       = TEMPLATES_DIR / 'model_diploma.pptx'
OUT_MEDICI     = OUTPUT_DIR    / 'Diplome_medici.pptx'
OUT_ASISTENTI  = OUTPUT_DIR    / 'Diplome_asistenti.pptx'
OUT_DIVERSE    = OUTPUT_DIR    / 'Diplome_diverse.pptx'


def _xlsx_rows(path):
    """
    Generator de dict-uri (header lowercase -> valoare) pentru fiecare rand
    de date dintr-un XLSX. Asteapta randul 1 = headere.
    """
    wb = openpyxl.load_workbook(path, data_only=True)
    ws = wb.active
    headers = [str(c.value).strip().lower() if c.value else '' for c in ws[1]]
    for row in ws.iter_rows(min_row=2, values_only=True):
        yield {h: (str(v).strip() if v else '') for h, v in zip(headers, row)}


def _load_event(path):
    """Citeste event.xlsx (primul rand de date: data, locatie)."""
    row = next(_xlsx_rows(path))
    return row['data'], row['locatie']


_data, _locatie   = _load_event(EVENT_XLSX)
EVENT_PERIOD_TEXT = f'desfășurat în perioada {_data},'
EVENT_CITY_TEXT   = _locatie

# Pozitii shape-uri in template
NAME_SHAPE_IDX    = 0
DETAILS_SHAPE_IDX = 2

# Dimensiuni font pentru Shape 0 (nume)
NAME_MAX_PT   = 36
NAME_MIN_PT   = 20
BOX_W_PT      = (9963000 - 116050 * 2) / 12700   # ≈ 766 pt
CALIBRI_RATIO = 0.55


# ─────────────────────────────────────────────────────────────────
# CITIRE CSV
# ─────────────────────────────────────────────────────────────────

def load_medici():
    """Returneaza lista de nume medici (unice, in ordinea aparitiei), inclusiv cei fara cabinet."""
    seen = set()
    out  = []
    for row in _xlsx_rows(MEDICI_XLSX):
        nume = row.get('nume', '')
        if nume and nume not in seen:
            seen.add(nume)
            out.append(nume)
    return out


def load_asistenti():
    """Returneaza lista de nume asistente."""
    out = []
    for row in _xlsx_rows(ASISTENTI_XLSX):
        nume = row.get('nume', '')
        if nume:
            out.append(nume)
    return out


def load_diversi():
    """
    Returneaza lista de "{titulatura} {nume}" pentru diplome diverse.
    Daca titulatura lipseste, foloseste doar numele.
    """
    out = []
    for row in _xlsx_rows(DIVERSE_XLSX):
        nume        = row.get('nume', '')
        titulatura  = row.get('titulatura', '')
        if not nume:
            continue
        out.append(f'{titulatura} {nume}'.strip())
    return out


# ─────────────────────────────────────────────────────────────────
# MANIPULARE SHAPE-URI IN SLIDE
# ─────────────────────────────────────────────────────────────────

def set_name_on_slide(slide, name: str):
    """Inlocuieste numele in Shape 0, cu font auto-scalat."""
    sh   = slide.shapes[NAME_SHAPE_IDX]
    para = sh._element.find(f'.//{{{NS_A}}}p')
    if para is None:
        return
    sz_pt = auto_font_pt(name, NAME_MAX_PT, NAME_MIN_PT, BOX_W_PT, CALIBRI_PATH, CALIBRI_RATIO)
    set_paragraph_single_run(para, name, sz_hundredths=sz_pt * 100)


def set_event_details_on_slide(slide):
    """Updateaza perioada (paragraful 0) si orasul (paragraful 1) din Shape 2."""
    if len(slide.shapes) <= DETAILS_SHAPE_IDX:
        return
    sh         = slide.shapes[DETAILS_SHAPE_IDX]
    paragraphs = sh._element.findall(f'.//{{{NS_A}}}p')
    if len(paragraphs) >= 1:
        set_paragraph_single_run(paragraphs[0], EVENT_PERIOD_TEXT)
    if len(paragraphs) >= 2:
        set_paragraph_single_run(paragraphs[1], EVENT_CITY_TEXT)


# ─────────────────────────────────────────────────────────────────
# GENERARE
# ─────────────────────────────────────────────────────────────────

def generate_pptx(names: list, output_path):
    prs = Presentation(str(TEMPLATE))
    template_slide = prs.slides[0]

    # Updateaza detaliile evenimentului INAINTE de clonare (asa apar pe toate slide-urile)
    set_event_details_on_slide(template_slide)

    # Slide 1: modifica template-ul direct
    set_name_on_slide(prs.slides[0], names[0])
    print(f'  Slide   1: {names[0]}')

    # Restul: clone din template
    for idx, name in enumerate(names[1:], start=2):
        new_slide = clone_slide(prs, template_slide)
        set_name_on_slide(new_slide, name)
        print(f'  Slide {idx:3d}: {name}')

    output_path.parent.mkdir(exist_ok=True)
    prs.save(str(output_path))
    print(f'\nSalvat: {output_path}  ({len(names)} slide-uri)')


def main():
    print(f'Font Calibri: {CALIBRI_PATH or "NE-GASIT (folosesc aproximare)"}')
    print(f'Eveniment: {EVENT_PERIOD_TEXT}  {EVENT_CITY_TEXT}\n')

    print('=== DIPLOME MEDICI ===')
    medici = load_medici()
    print(f'  Medici unici: {len(medici)}')
    generate_pptx(medici, OUT_MEDICI)

    print('\n=== DIPLOME ASISTENTI ===')
    asistenti = load_asistenti()
    print(f'  Asistente: {len(asistenti)}')
    generate_pptx(asistenti, OUT_ASISTENTI)

    if DIVERSE_XLSX.exists():
        print('\n=== DIPLOME DIVERSE ===')
        diversi = load_diversi()
        if diversi:
            print(f'  Persoane: {len(diversi)}')
            generate_pptx(diversi, OUT_DIVERSE)
        else:
            print('  (lista goala — sar peste)')
    else:
        print(f'\n(skip diplome diverse: nu exista {DIVERSE_XLSX.name})')

    print('\nGata.')


if __name__ == '__main__':
    main()
